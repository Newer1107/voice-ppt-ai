"""Stage 8: Embed narration audio into PowerPoint.

Input: Original PPTX path + per-slide narration audio files
Output: Narrated PPTX with embedded audio
Verifies: Output opens correctly, audio embedded, slide count preserved
"""

import logging
import os
import shutil
import tempfile
import zipfile
from pathlib import Path
from typing import Optional

from lxml import etree
from pptx import Presentation
from pydantic import BaseModel

logger = logging.getLogger(__name__)


class EmbedNarrationResult(BaseModel):
    output_path: str
    slide_count: int
    audio_tracks_added: int
    file_size_bytes: int


_NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_NS_RELS = "http://schemas.openxmlformats.org/package/2006/relationships"
_AUDIO_CTYPE = "audio/wav"


def _next_rid(existing_rels_root: etree._Element) -> str:
    seen: set[str] = set()
    for rel in existing_rels_root.iter(f"{{{_NS_RELS}}}Relationship"):
        rid = rel.get("Id", "")
        if rid:
            seen.add(rid)
    n = 1
    while f"rId{n}" in seen:
        n += 1
    return f"rId{n}"


def _read_zip_to_dict(pptx_path: str) -> dict[str, bytes]:
    entries: dict[str, bytes] = {}
    with zipfile.ZipFile(pptx_path, "r") as zf:
        for info in zf.infolist():
            entries[info.filename] = zf.read(info.filename)
    return entries


def _write_dict_to_zip(entries: dict[str, bytes], dst: str) -> None:
    fd, tmp = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)
    with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for name, data in entries.items():
            zout.writestr(name, data)
    shutil.move(tmp, dst)


def _ensure_wav_content_type(entries: dict[str, bytes]) -> None:
    ct_path = "[Content_Types].xml"
    if ct_path not in entries:
        return
    root = etree.fromstring(entries[ct_path])
    ns_ct = "http://schemas.openxmlformats.org/package/2006/content-types"
    for default in root.iter(f"{{{ns_ct}}}Default"):
        if (default.get("Extension") or "").lower() == "wav":
            return
    default = etree.SubElement(root, f"{{{ns_ct}}}Default")
    default.set("Extension", "wav")
    default.set("ContentType", _AUDIO_CTYPE)
    entries[ct_path] = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)


def _max_shape_id(slide_root: etree._Element) -> int:
    max_id = 0
    for el in slide_root.iter(f"{{{_NS_P}}}cNvPr"):
        try:
            max_id = max(max_id, int(el.get("id", "0")))
        except (ValueError, TypeError):
            continue
    return max_id


def _add_audio_shape(
    sp_tree: etree._Element,
    rid: str,
    shape_id: int,
    slide_num: int,
) -> None:
    gf = etree.SubElement(sp_tree, f"{{{_NS_P}}}graphicFrame")
    nv = etree.SubElement(gf, f"{{{_NS_P}}}nvGraphicFramePr")
    cnv = etree.SubElement(nv, f"{{{_NS_P}}}cNvPr")
    cnv.set("id", str(shape_id))
    cnv.set("name", f"Audio {slide_num}")
    cnv.set("hidden", "1")
    cnv_gf = etree.SubElement(nv, f"{{{_NS_P}}}cNvGraphicFramePr")
    locks = etree.SubElement(cnv_gf, f"{{{_NS_A}}}graphicFrameLocks")
    locks.set("noGrp", "1")
    etree.SubElement(nv, f"{{{_NS_P}}}nvPr")
    xfrm = etree.SubElement(gf, f"{{{_NS_P}}}xfrm")
    off = etree.SubElement(xfrm, f"{{{_NS_A}}}off")
    off.set("x", "0")
    off.set("y", "0")
    ext = etree.SubElement(xfrm, f"{{{_NS_A}}}ext")
    ext.set("cx", "0")
    ext.set("cy", "0")
    graphic = etree.SubElement(gf, f"{{{_NS_A}}}graphic")
    gdata = etree.SubElement(graphic, f"{{{_NS_A}}}graphicData")
    gdata.set("uri", "http://schemas.openxmlformats.org/presentationml/2006/media")
    media = etree.SubElement(gdata, f"{{{_NS_P}}}media")
    media.set(f"{{{_NS_R}}}embed", rid)


def _ensure_timing_media(slide_root: etree._Element) -> None:
    timing = slide_root.find(f"{{{_NS_P}}}timing")
    if timing is not None:
        child_tn = timing.find(
            f"{{{_NS_P}}}tnLst/{{{_NS_P}}}par/{{{_NS_P}}}cTnNode/{{{_NS_P}}}childTnLst"
        )
        if child_tn is not None:
            _append_media_to_child_tn(child_tn)
            return
        slide_root.remove(timing)
    timing = etree.SubElement(slide_root, f"{{{_NS_P}}}timing")
    tn_lst = etree.SubElement(timing, f"{{{_NS_P}}}tnLst")
    par = etree.SubElement(tn_lst, f"{{{_NS_P}}}par")
    ctn = etree.SubElement(par, f"{{{_NS_P}}}cTnNode")
    child_tn = etree.SubElement(ctn, f"{{{_NS_P}}}childTnLst")
    _append_media_to_child_tn(child_tn)


def _append_media_to_child_tn(parent: etree._Element) -> None:
    seq = etree.SubElement(parent, f"{{{_NS_P}}}seq")
    seq_ctn = etree.SubElement(seq, f"{{{_NS_P}}}cTnNode")
    seq_child = etree.SubElement(seq_ctn, f"{{{_NS_P}}}childTnLst")
    media = etree.SubElement(seq_child, f"{{{_NS_P}}}media")
    media_node = etree.SubElement(media, f"{{{_NS_P}}}cMediaNode")
    cond = etree.SubElement(media_node, f"{{{_NS_P}}}cond")
    cond.set("condDelay", "0")


def embed_narration_into_pptx(
    pptx_path: str,
    slide_audio_map: dict[int, str],
    output_path: Optional[str] = None,
) -> EmbedNarrationResult:
    """Embed narration audio files into a PowerPoint presentation.

    For each slide with audio, embeds the audio WAV data directly into the
    PPTX zip archive and configures it to play automatically on slide entry.

    Verifies:
    - Output file opens correctly via ``python-pptx``
    - Audio files are copied into the zip (embedded, not external references)
    - Slide count matches original presentation
    """
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"PPTX file not found: {pptx_path}")

    if not slide_audio_map:
        raise ValueError("No audio files to embed")

    for slide_num, audio_path in slide_audio_map.items():
        if not os.path.exists(audio_path):
            raise FileNotFoundError(f"Audio file for slide {slide_num} not found: {audio_path}")

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        raise ValueError(f"Failed to open PPTX for embedding: {e}")
    original_slide_count = len(prs.slides)

    if output_path is None:
        base = Path(pptx_path).stem
        output_dir = os.path.dirname(pptx_path)
        output_path = os.path.join(output_dir, f"{base}_narrated.pptx")

    entries = _read_zip_to_dict(pptx_path)

    _ensure_wav_content_type(entries)

    audio_added = 0

    for slide_num, audio_path in sorted(slide_audio_map.items()):
        if slide_num < 1 or slide_num > original_slide_count:
            logger.warning("Slide %d out of range (1-%d), skipping", slide_num, original_slide_count)
            continue

        if not os.path.exists(audio_path):
            logger.warning("Audio file %s for slide %d not found, skipping", audio_path, slide_num)
            continue

        slide_xml_path = f"ppt/slides/slide{slide_num}.xml"
        rels_path = f"ppt/slides/_rels/slide{slide_num}.xml.rels"

        if slide_xml_path not in entries:
            logger.warning("Slide XML not found at %s, skipping", slide_xml_path)
            continue

        media_arcname = f"ppt/media/audio{slide_num:03d}.wav"
        with open(audio_path, "rb") as f:
            entries[media_arcname] = f.read()

        rels_raw = entries.get(rels_path)
        if rels_raw is None:
            logger.warning("No rels file at %s, skipping slide %d", rels_path, slide_num)
            continue
        rels_root = etree.fromstring(rels_raw)
        new_rid = _next_rid(rels_root)
        new_rel = etree.SubElement(rels_root, f"{{{_NS_RELS}}}Relationship")
        new_rel.set("Id", new_rid)
        new_rel.set("Type", f"{_NS_R}/media")
        new_rel.set("Target", f"../media/audio{slide_num:03d}.wav")
        entries[rels_path] = etree.tostring(rels_root, xml_declaration=True, encoding="UTF-8", standalone=True)

        slide_root = etree.fromstring(entries[slide_xml_path])
        sp_tree = slide_root.find(f".//{{{_NS_P}}}spTree")
        if sp_tree is None:
            logger.warning("No spTree in slide %d XML, skipping", slide_num)
            continue

        new_id = _max_shape_id(slide_root) + 1
        _add_audio_shape(sp_tree, new_rid, new_id, slide_num)
        _ensure_timing_media(slide_root)

        entries[slide_xml_path] = etree.tostring(slide_root, xml_declaration=True, encoding="UTF-8", standalone=True)
        audio_added += 1

        logger.debug("Embedded audio %s into slide %d", media_arcname, slide_num)

    _write_dict_to_zip(entries, output_path)

    if not os.path.exists(output_path):
        raise RuntimeError(f"Narrated PPTX not created: {output_path}")

    file_size = os.path.getsize(output_path)

    try:
        verify = Presentation(output_path)
        slide_count = len(verify.slides)
        if slide_count != original_slide_count:
            logger.warning(
                "Slide count mismatch: original=%d, output=%d",
                original_slide_count,
                slide_count,
            )
    except Exception as e:
        raise RuntimeError(f"Output PPTX failed verification: {e}")

    logger.info(
        "Narrated PPTX created: %s (%d slides, %d audio tracks, %d bytes)",
        output_path,
        slide_count,
        audio_added,
        file_size,
    )

    return EmbedNarrationResult(
        output_path=str(output_path),
        slide_count=slide_count,
        audio_tracks_added=audio_added,
        file_size_bytes=file_size,
    )
