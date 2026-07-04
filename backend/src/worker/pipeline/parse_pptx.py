"""Stage 3: Parse PowerPoint file to extract structured slide content.

Input: PPTX file path
Output: List of structured slide data
Verifies: Slides extracted, count matches presentation
"""

import logging
import os
from typing import Optional

from pptx import Presentation
from pydantic import BaseModel

logger = logging.getLogger(__name__)


class ParsedSlide(BaseModel):
    """Structured content of a single slide."""
    slide_number: int
    raw_text: str
    notes: Optional[str] = None
    slide_layout: Optional[str] = None
    image_count: int = 0


class PptxParseResult(BaseModel):
    """Result of PPTX parsing."""
    slides: list[ParsedSlide]
    total_slides: int
    slide_width: int
    slide_height: int


def parse_pptx(pptx_path: str) -> PptxParseResult:
    """Parse a PowerPoint file and extract all slide content.

    Extracts:
    - All text from all shapes on each slide
    - Speaker notes
    - Layout information
    - Image count per slide

    Raises FileNotFoundError if file missing, ValueError if invalid.
    """
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"PPTX file not found: {pptx_path}")

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        raise ValueError(f"Failed to open PPTX: {e}")

    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        all_text = []
        image_count = 0

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        all_text.append(text)
            if shape.shape_type == 13:  # Picture
                image_count += 1

        # Extract speaker notes
        notes_text = None
        if slide.has_notes_slide:
            try:
                notes_text = slide.notes_slide.notes_text_frame.text.strip() or None
            except Exception:
                pass

        slides.append(ParsedSlide(
            slide_number=i,
            raw_text="\n".join(all_text),
            notes=notes_text,
            slide_layout=slide.slide_layout.name if slide.slide_layout else None,
            image_count=image_count,
        ))

    if not slides:
        raise ValueError(f"No slides found in presentation: {pptx_path}")

    logger.info(
        "Parsed PPTX: %s -> %d slides (%.1f x %.1f)",
        pptx_path, len(slides),
        prs.slide_width, prs.slide_height,
    )

    return PptxParseResult(
        slides=slides,
        total_slides=len(slides),
        slide_width=prs.slide_width,
        slide_height=prs.slide_height,
    )
